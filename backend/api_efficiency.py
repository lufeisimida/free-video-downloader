"""APIs for automated learning, retrieval, media, reminders and quality control."""

from __future__ import annotations

import asyncio
import io
import ipaddress
import json
import mimetypes
import os
import socket
import tempfile
import time
from datetime import date, datetime, timedelta, timezone
from html.parser import HTMLParser
from pathlib import Path
from urllib.parse import urlparse
from uuid import uuid4

import httpx
from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, RedirectResponse
from pydantic import BaseModel, Field

from auth import get_current_user
from content_store import content_store
from database import (
    get_library_video,
    get_video_folder,
    list_flashcards,
    list_folder_library_videos,
    list_knowledge_mastery,
    list_question_bank,
    list_video_progress,
    list_wrong_questions,
)
from learning_efficiency import (
    assess_question_quality,
    create_material,
    create_processing_job,
    deduplicate_knowledge_points,
    delete_material,
    get_material,
    get_processing_job,
    get_reminder_preferences,
    index_source,
    list_knowledge_aliases,
    list_materials,
    list_pending_browser_reminders,
    list_processing_jobs,
    list_queued_processing_jobs,
    list_question_quality,
    mark_reminder_notified,
    save_question_quality,
    save_reminder_preferences,
    semantic_search,
    update_material_status,
    subtitle_chunks,
    update_job,
    update_job_item,
    usage_summary,
)
from local_media import is_local_media_url, read_metadata, resolve_local_path


router = APIRouter(prefix="/api/library/efficiency", tags=["学习效率"])
_pipeline_tasks: dict[int, asyncio.Task] = {}
_material_tasks: dict[int, asyncio.Task] = {}
_media_tickets: dict[str, dict] = {}
_pipeline_semaphore: asyncio.Semaphore | None = None


class PipelineRequest(BaseModel):
    force: bool = False
    generate_questions: bool = True
    priority: int = Field(default=50, ge=1, le=100)


class SemanticRequest(BaseModel):
    query: str
    limit: int = Field(default=8, ge=1, le=30)


class CourseQuestionRequest(SemanticRequest):
    pass


class WebMaterialRequest(BaseModel):
    url: str


class ReminderRequest(BaseModel):
    enabled: bool = False
    reminder_time: str = "20:00"
    browser_enabled: bool = True
    email_enabled: bool = False
    wecom_enabled: bool = False


class ReminderNotifiedRequest(BaseModel):
    delivery_ids: list[int] = Field(default_factory=list)


class PlaybackRequest(BaseModel):
    format_id: str = "best"


class MaterialProcessRequest(BaseModel):
    force: bool = False
    question_count: int = Field(default=8, ge=2, le=20)


def _ensure_folder(user_id: int, folder_id: int) -> dict:
    folder = get_video_folder(user_id, folder_id)
    if not folder:
        raise HTTPException(status_code=404, detail="课程目录不存在")
    return folder


def _has_subtitle(value: object) -> bool:
    return isinstance(value, dict) and bool(str(value.get("full_text") or "").strip())


async def _run_pipeline(job_id: int, user_id: int, folder_id: int,
                        generate_questions: bool, force: bool) -> None:
    global _pipeline_semaphore
    if _pipeline_semaphore is None:
        _pipeline_semaphore = asyncio.Semaphore(max(1, int(os.getenv("PIPELINE_MAX_CONCURRENCY", "2"))))
    async with _pipeline_semaphore:
        await _run_pipeline_locked(job_id, user_id, folder_id, generate_questions, force)


async def _run_pipeline_locked(job_id: int, user_id: int, folder_id: int,
                               generate_questions: bool, force: bool) -> None:
    job = get_processing_job(user_id, job_id)
    if not job:
        return
    update_job(job_id, status="running", stage="parsing", progress=1, error="")
    completed = 0
    failed = 0
    loop = asyncio.get_running_loop()
    for item in job["items"]:
        video_id = item["video_id"]
        video = get_library_video(user_id, video_id)
        if not video:
            failed += 1
            update_job_item(job_id, video_id, status="failed", stage="parsing", progress=0, error="视频记录不存在")
            continue
        try:
            update_job_item(job_id, video_id, status="running", stage="subtitle", progress=20)
            update_job(job_id, stage="subtitle")
            saved = content_store.load(video["url"])
            subtitle = saved.get("subtitle")
            if force or not _has_subtitle(subtitle):
                from api_summarize import _get_extractor, _get_transcriber

                subtitle = await loop.run_in_executor(None, _get_extractor().extract, video["url"])
                if not _has_subtitle(subtitle):
                    subtitle = await loop.run_in_executor(None, _get_transcriber().transcribe_url, video["url"])
                content_store.update(video["url"], title=video["title"], subtitle=subtitle)
            if not _has_subtitle(subtitle):
                raise ValueError("没有取得可用字幕")

            update_job_item(job_id, video_id, status="running", stage="indexing", progress=55)
            update_job(job_id, stage="indexing")
            chunks = subtitle_chunks(subtitle)
            index_source(
                user_id, folder_id, content=subtitle.get("full_text", ""),
                source_title=video["title"], source_url=video["url"],
                source_type="video", video_id=video_id, timed_chunks=chunks,
            )

            update_job_item(job_id, video_id, status="running", stage="knowledge", progress=75)
            update_job(job_id, stage="knowledge")
            deduplicate_knowledge_points(user_id, folder_id)

            if generate_questions:
                update_job_item(job_id, video_id, status="running", stage="question_bank", progress=88)
                update_job(job_id, stage="question_bank")
                existing = list_question_bank(user_id, folder_id, limit=1)
                if force or not existing:
                    from api_library import _persist_questions, _prepare_question
                    from api_summarize import _get_summarizer

                    spec = {"type": "short_answer", "label": "简答题", "count": 2, "points": 5}
                    payload = await loop.run_in_executor(
                        None,
                        _get_summarizer("quiz").generate_quiz_batch,
                        subtitle.get("full_text", "")[:12000], "zh", spec, 1, [],
                    )
                    questions = [
                        _prepare_question(question, [(video, subtitle.get("full_text", ""))])
                        for question in payload.get("questions", [])
                    ]
                    _persist_questions(user_id, folder_id, questions, subtitle.get("full_text", ""))

            completed += 1
            update_job_item(job_id, video_id, status="completed", stage="completed", progress=100)
        except Exception as exc:
            failed += 1
            update_job_item(job_id, video_id, status="failed", stage="failed", progress=100, error=str(exc))
        progress = round((completed + failed) * 100 / max(job["total_items"], 1), 1)
        update_job(job_id, progress=progress, completed_items=completed)
    status = "failed" if failed and not completed else "completed"
    message = f"{failed} 个视频处理失败，可点击重试" if failed else ""
    update_job(job_id, status=status, stage=status, progress=100,
               completed_items=completed, error=message)
    _pipeline_tasks.pop(job_id, None)


def _schedule_pipeline(job: dict, user_id: int, folder_id: int,
                       generate_questions: bool, force: bool) -> None:
    if job["id"] in _pipeline_tasks and not _pipeline_tasks[job["id"]].done():
        return
    task = asyncio.create_task(_run_pipeline(job["id"], user_id, folder_id, generate_questions, force))
    _pipeline_tasks[job["id"]] = task


def resume_queued_pipelines() -> int:
    jobs = list_queued_processing_jobs()
    for job in jobs:
        _schedule_pipeline(
            job, job["user_id"], job["folder_id"], bool(job.get("generate_questions", 1)),
            bool(job.get("force", 0)),
        )
    return len(jobs)


@router.post("/folders/{folder_id}/pipeline")
async def start_pipeline(folder_id: int, req: PipelineRequest,
                         user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    videos = list_folder_library_videos(user["id"], folder_id)
    if not videos:
        raise HTTPException(status_code=400, detail="课程目录内没有视频")
    job = create_processing_job(
        user["id"], folder_id, [item["id"] for item in videos], req.priority,
        req.generate_questions, req.force,
    )
    _schedule_pipeline(job, user["id"], folder_id, req.generate_questions, req.force)
    return {"success": True, "data": job}


@router.get("/folders/{folder_id}/pipeline")
async def processing_jobs(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    return {"success": True, "data": list_processing_jobs(user["id"], folder_id)}


@router.post("/jobs/{job_id}/retry")
async def retry_pipeline(job_id: int, user: dict = Depends(get_current_user)):
    job = get_processing_job(user["id"], job_id)
    if not job:
        raise HTTPException(status_code=404, detail="处理任务不存在")
    if job_id in _pipeline_tasks and not _pipeline_tasks[job_id].done():
        raise HTTPException(status_code=409, detail="任务仍在运行")
    update_job(job_id, status="queued", stage="queued", progress=0, completed_items=0, error="")
    _schedule_pipeline(job, user["id"], job["folder_id"], True, False)
    return {"success": True, "data": get_processing_job(user["id"], job_id)}


@router.post("/folders/{folder_id}/semantic-search")
async def course_semantic_search(folder_id: int, req: SemanticRequest,
                                 user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    query = req.query.strip()
    if len(query) < 2:
        raise HTTPException(status_code=400, detail="检索内容至少需要 2 个字符")
    return {"success": True, "data": {
        "query": query, "results": semantic_search(user["id"], folder_id, query, req.limit),
    }}


@router.post("/folders/{folder_id}/ask")
async def ask_course(folder_id: int, req: CourseQuestionRequest,
                     user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    question = req.query.strip()
    sources = semantic_search(user["id"], folder_id, question, req.limit)
    if not sources:
        raise HTTPException(status_code=400, detail="课程尚未建立语义索引，请先运行自动处理")
    evidence = "\n\n".join(
        f"[证据 {index + 1}] {item['source_title']} {item['start_seconds']:.0f}秒\n{item['content']}"
        for index, item in enumerate(sources)
    )
    prompt = (
        "只根据以下课程证据回答问题。证据不足时明确说不知道，不要补充课程外事实。"
        "回答中使用[证据1]格式标注依据。\n\n" + evidence + "\n\n问题：" + question
    )
    try:
        from api_summarize import _get_summarizer

        summarizer = _get_summarizer("course_qa")
        answer = await asyncio.get_running_loop().run_in_executor(
            None,
            lambda: summarizer.provider.complete_chat(
                [{"role": "system", "content": "你是严格依据课程材料回答的学习助教。"},
                 {"role": "user", "content": prompt}],
                temperature=0.2, max_tokens=1400,
            ),
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"课程问答失败: {exc}") from exc
    citations = [{
        "source_title": item["source_title"], "source_url": item["source_url"],
        "source_type": item["source_type"], "time_seconds": item["start_seconds"],
        "quote": item["content"][:260], "score": item["score"],
    } for item in sources]
    return {"success": True, "data": {"answer": answer, "citations": citations}}


@router.get("/folders/{folder_id}/continuous")
async def continuous_learning(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    queue: list[dict] = []
    for card in list_flashcards(user["id"], folder_id, due_only=True)[:12]:
        queue.append({"key": f"card:{card['id']}", "type": "flashcard", "title": card["front"], "data": card})
    for item in list_wrong_questions(user["id"], folder_id, due_only=True)[:8]:
        queue.append({"key": f"wrong:{item['id']}", "type": "wrong", "title": item["question"].get("question", "错题复习"), "data": item})
    video = next((item for item in list_video_progress(user["id"], folder_id) if item["status"] != "completed"), None)
    if video:
        queue.append({"key": f"video:{video['video_id']}", "type": "video",
                      "title": f"继续学习：{video['title']}", "data": video})
    mastery = list_knowledge_mastery(user["id"], folder_id)
    quiz_data = {"mode": "adaptive", "phase": "practice"}
    if mastery:
        quiz_data["knowledge_point"] = mastery[0]["knowledge_point"]
    queue.append({"key": "quiz:adaptive", "type": "quiz", "title": "完成今日小测", "data": quiz_data})
    return {"success": True, "data": {"items": queue, "total": len(queue)}}


@router.get("/reminders")
async def reminder_settings(user: dict = Depends(get_current_user)):
    return {"success": True, "data": get_reminder_preferences(user["id"])}


@router.put("/reminders")
async def update_reminders(req: ReminderRequest, user: dict = Depends(get_current_user)):
    try:
        datetime.strptime(req.reminder_time, "%H:%M")
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="提醒时间必须是 HH:MM") from exc
    return {"success": True, "data": save_reminder_preferences(
        user["id"], req.enabled, req.reminder_time, req.browser_enabled, req.email_enabled,
        req.wecom_enabled,
    )}


@router.get("/reminders/pending")
async def pending_reminders(user: dict = Depends(get_current_user)):
    return {"success": True, "data": list_pending_browser_reminders(user["id"])}


@router.post("/reminders/notified")
async def reminder_was_shown(req: ReminderNotifiedRequest,
                             user: dict = Depends(get_current_user)):
    from learning_efficiency import acknowledge_browser_reminders

    acknowledge_browser_reminders(user["id"], req.delivery_ids)
    if not req.delivery_ids:
        mark_reminder_notified(user["id"])
    return {"success": True}


class _TextHTMLParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []
        self.ignored = 0

    def handle_starttag(self, tag, attrs):
        if tag in {"script", "style", "noscript"}:
            self.ignored += 1

    def handle_endtag(self, tag):
        if tag in {"script", "style", "noscript"} and self.ignored:
            self.ignored -= 1

    def handle_data(self, data):
        if not self.ignored and data.strip():
            self.parts.append(data.strip())


def _validate_public_url(url: str) -> None:
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise HTTPException(status_code=400, detail="仅支持 http/https 网页地址")
    try:
        addresses = socket.getaddrinfo(parsed.hostname, parsed.port or (443 if parsed.scheme == "https" else 80))
    except socket.gaierror as exc:
        raise HTTPException(status_code=400, detail="网页域名无法解析") from exc
    for result in addresses:
        ip = ipaddress.ip_address(result[4][0])
        if not ip.is_global:
            raise HTTPException(status_code=400, detail="不允许导入内网或本机地址")


@router.post("/folders/{folder_id}/materials/web")
async def import_web_material(folder_id: int, req: WebMaterialRequest,
                              user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    _validate_public_url(req.url)
    try:
        async with httpx.AsyncClient(timeout=20, follow_redirects=False) as client:
            response = await client.get(req.url, headers={"User-Agent": "SaveAnyLearning/1.0"})
            response.raise_for_status()
        parser = _TextHTMLParser()
        parser.feed(response.text[:5_000_000])
        content = "\n".join(parser.parts)
        if len(content) < 20:
            raise ValueError("网页没有可提取的正文")
        name = urlparse(req.url).hostname or "网页资料"
        material = create_material(user["id"], folder_id, name, "web", content, req.url)
        index_source(user["id"], folder_id, content=content, source_title=name,
                     source_url=req.url, source_type="web", material_id=material["id"])
        return {"success": True, "data": material}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"网页导入失败: {exc}") from exc


def _extract_uploaded_text(filename: str, raw: bytes) -> tuple[str, str]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise ValueError("缺少 pypdf，请重新安装 requirements.txt") from exc
        reader = PdfReader(io.BytesIO(raw))
        return "\n".join(page.extract_text() or "" for page in reader.pages), "pdf"
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ValueError("缺少 python-pptx，请重新安装 requirements.txt") from exc
        deck = Presentation(io.BytesIO(raw))
        text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
        return text, "pptx"
    if suffix in {".txt", ".md", ".csv", ".srt", ".vtt"}:
        return raw.decode("utf-8", errors="replace"), suffix.removeprefix(".")
    raise ValueError("仅支持 PDF、PPTX、TXT、Markdown、CSV、SRT 和 VTT")


@router.post("/folders/{folder_id}/materials/upload")
async def upload_material(folder_id: int, file: UploadFile = File(...),
                          user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    raw = await file.read(30 * 1024 * 1024 + 1)
    await file.close()
    if len(raw) > 30 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="单个资料不能超过 30MB")
    name = Path(file.filename or "资料").name
    try:
        content, material_type = _extract_uploaded_text(name, raw)
        content = content.strip()
        if len(content) < 10:
            raise ValueError("资料中没有提取到足够文字，扫描版 PDF 请先 OCR")
        material = create_material(user["id"], folder_id, name, material_type, content)
        index_source(user["id"], folder_id, content=content, source_title=name,
                     source_type=material_type, material_id=material["id"])
        return {"success": True, "data": material}
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


async def _process_material(material: dict, question_count: int) -> None:
    material_id = int(material["id"])
    update_material_status(material_id, "processing", "")
    try:
        from api_library import _persist_questions
        from api_summarize import _get_summarizer

        content = str(material.get("content") or "")
        spec = {
            "type": "short_answer", "label": "资料简答题",
            "count": question_count, "points": 5, "max_tokens": 4096,
        }
        payload = await asyncio.to_thread(
            _get_summarizer("quiz").generate_quiz_batch, content[:20000], "zh", spec, 1, [],
        )
        source = {
            "id": None, "title": material["name"],
            "url": material.get("source_url") or f"material://{material_id}",
        }
        questions = []
        for question in payload.get("questions", []):
            item = dict(question)
            item["source_video_title"] = source["title"]
            item["source_video_url"] = source["url"]
            item["evidence_time_seconds"] = 0
            item.setdefault("knowledge_point", material["name"][:80])
            item.setdefault("evidence_quote", content[:300])
            from api_library import _prepare_question

            questions.append(_prepare_question(item, [(source, content)]))
        saved = _persist_questions(material["user_id"], material["folder_id"], questions, content)
        if not saved:
            raise ValueError("生成的题目均未通过质量校验")
        deduplicate_knowledge_points(material["user_id"], material["folder_id"])
        update_material_status(material_id, "ready", "")
    except Exception as exc:
        update_material_status(material_id, "failed", str(exc))
    finally:
        _material_tasks.pop(material_id, None)


@router.post("/materials/{material_id}/process")
async def process_material(material_id: int, req: MaterialProcessRequest,
                           user: dict = Depends(get_current_user)):
    material = get_material(user["id"], material_id)
    if not material:
        raise HTTPException(status_code=404, detail="资料不存在")
    if material_id in _material_tasks and not _material_tasks[material_id].done():
        raise HTTPException(status_code=409, detail="资料正在处理中")
    task = asyncio.create_task(_process_material(material, req.question_count))
    _material_tasks[material_id] = task
    return {"success": True, "data": {**material, "status": "processing"}}


@router.get("/folders/{folder_id}/materials")
async def materials(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    return {"success": True, "data": list_materials(user["id"], folder_id)}


@router.delete("/materials/{material_id}")
async def remove_material(material_id: int, user: dict = Depends(get_current_user)):
    if not delete_material(user["id"], material_id):
        raise HTTPException(status_code=404, detail="资料不存在")
    return {"success": True}


@router.post("/folders/{folder_id}/knowledge/deduplicate")
async def deduplicate_points(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    return {"success": True, "data": deduplicate_knowledge_points(user["id"], folder_id)}


@router.get("/folders/{folder_id}/knowledge/aliases")
async def knowledge_aliases(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    return {"success": True, "data": list_knowledge_aliases(user["id"], folder_id)}


@router.post("/folders/{folder_id}/quality/check")
async def check_question_quality(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    questions = list_question_bank(user["id"], folder_id, limit=1000)
    rejected = 0
    for index, item in enumerate(questions):
        quality = assess_question_quality(item["question"], questions[:index] + questions[index + 1:])
        if quality["status"] == "rejected":
            from database import delete_question_bank_item

            rejected += int(delete_question_bank_item(user["id"], item["id"]))
            continue
        save_question_quality(item["id"], quality)
    return {"success": True, "data": list_question_quality(user["id"], folder_id),
            "meta": {"rejected_removed": rejected}}


@router.get("/folders/{folder_id}/quality")
async def question_quality(folder_id: int, user: dict = Depends(get_current_user)):
    _ensure_folder(user["id"], folder_id)
    return {"success": True, "data": list_question_quality(user["id"], folder_id)}


@router.get("/usage")
async def model_usage(days: int = 30, user: dict = Depends(get_current_user)):
    return {"success": True, "data": usage_summary(user["id"], days)}


@router.post("/videos/{video_id}/playback")
async def create_playback(video_id: int, req: PlaybackRequest,
                          user: dict = Depends(get_current_user)):
    video = get_library_video(user["id"], video_id)
    if not video:
        raise HTTPException(status_code=404, detail="视频记录不存在")
    if is_local_media_url(video["url"]):
        ticket = uuid4().hex
        _media_tickets[ticket] = {
            "user_id": user["id"], "video_id": video_id,
            "expires_at": datetime.now(timezone.utc) + timedelta(minutes=15),
        }
        return {"success": True, "data": {
            "url": f"/api/library/efficiency/media/{video_id}?ticket={ticket}",
            "expires_in": 900, "local": True,
        }}
    try:
        from downloader import VideoDownloader

        result = await asyncio.get_running_loop().run_in_executor(
            None, VideoDownloader().get_direct_url, video["url"], req.format_id,
        )
        url = result.get("url") if isinstance(result, dict) else result
        return {"success": True, "data": {"url": url, "expires_in": 600, "local": False}}
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"获取播放地址失败: {exc}") from exc


@router.get("/media/{video_id}")
async def stream_local_media(video_id: int, ticket: str):
    entry = _media_tickets.get(ticket)
    if not entry or entry["video_id"] != video_id or entry["expires_at"] < datetime.now(timezone.utc):
        _media_tickets.pop(ticket, None)
        raise HTTPException(status_code=401, detail="播放凭证无效或已过期")
    video = get_library_video(entry["user_id"], video_id)
    path = resolve_local_path(video["url"]) if video else None
    if not path:
        raise HTTPException(status_code=404, detail="本地视频不存在")
    metadata = read_metadata(Path(path).name)
    media_type = metadata.get("content_type") or mimetypes.guess_type(path)[0] or "video/mp4"
    return FileResponse(path, media_type=media_type)


@router.post("/voice/recall")
async def voice_recall(
    audio: UploadFile = File(...),
    reference: str = Form(...),
    question: str = Form(""),
    user: dict = Depends(get_current_user),
):
    raw = await audio.read(20 * 1024 * 1024 + 1)
    suffix = Path(audio.filename or "recall.webm").suffix or ".webm"
    await audio.close()
    if len(raw) > 20 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="语音文件不能超过 20MB")
    temp_path = ""
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as output:
            output.write(raw)
            temp_path = output.name
        from api_summarize import _get_transcriber

        transcriber = _get_transcriber()

        def transcribe_file() -> str:
            model = transcriber._get_model()
            segments, _ = model.transcribe(temp_path, language=transcriber.language,
                                           beam_size=transcriber.beam_size, vad_filter=True)
            return " ".join(str(segment.text or "").strip() for segment in segments).strip()

        transcript = await asyncio.get_running_loop().run_in_executor(None, transcribe_file)
        if not transcript:
            raise ValueError("没有识别出有效语音")
        from learning_efficiency import tokenize

        expected = set(tokenize(reference))
        actual = set(tokenize(transcript))
        recall = len(expected & actual) / max(len(expected), 1)
        precision = len(expected & actual) / max(len(actual), 1)
        score = round((recall * 0.7 + precision * 0.3) * 100)
        missing = list(expected - actual)[:12]
        feedback = "复述准确，可以进入下一项。" if score >= 80 else "核心内容基本覆盖，补充遗漏点后再复述一次。" if score >= 60 else "遗漏较多，先查看参考答案再重新复述。"
        return {"success": True, "data": {
            "question": question, "transcript": transcript, "score": score,
            "feedback": feedback, "missing_keywords": missing,
        }}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"语音复述评分失败: {exc}") from exc
    finally:
        if temp_path:
            try:
                os.unlink(temp_path)
            except OSError:
                pass
